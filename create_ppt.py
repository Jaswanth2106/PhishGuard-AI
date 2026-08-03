import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Base paths
    image_dir = r"C:\Users\user\.gemini\antigravity\brain\7fde9697-e62a-4aee-9ca4-f1be1f732622\.user_uploaded"
    
    img_logo = os.path.join(image_dir, "media__1785767475143.png")
    img_analyse = os.path.join(image_dir, "media__1785767475168.png")
    img_gemini = os.path.join(image_dir, "media__1785767475175.png")
    img_chat = os.path.join(image_dir, "media__1785767475185.png")
    img_details = os.path.join(image_dir, "media__1785767475195.png")
    img_history = os.path.join(image_dir, "media__1785767475204.png")

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "PhishGuard AI"
    subtitle.text = "Intelligent Threat Detection & Explainable Cybersecurity\n\nBuilt entirely with OpenAI Codex"
    
    # Slide 2: The Core Problem & Solution (Logo)
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Project Overview"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(4))
    tf = txBox.text_frame
    tf.text = "The Problem:\n- Traditional filters block without explaining.\n- Users remain uneducated.\n- Image-based threats bypass scanners.\n\nThe Solution (PhishGuard AI):\n- Next.js & FastAPI Architecture.\n- Real-time ML Classification.\n- Gemini 1.5 Flash for Explainable AI (XAI)."
    
    if os.path.exists(img_logo):
        slide.shapes.add_picture(img_logo, Inches(5), Inches(1.5), height=Inches(4))

    # Slide 3: Real-Time Detection
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Feature: Real-Time Threat Detection"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.text = "Our Scikit-learn ML model processes incoming emails instantly, providing a mathematical confidence score."
    
    if os.path.exists(img_analyse):
        slide.shapes.add_picture(img_analyse, Inches(0.5), Inches(2.5), width=Inches(9))

    # Slide 4: Explainable AI & Technical Details
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Feature: Explainable AI (XAI) via Gemini"
    
    if os.path.exists(img_gemini):
        slide.shapes.add_picture(img_gemini, Inches(0.5), Inches(1.5), height=Inches(5))
        
    if os.path.exists(img_details):
        slide.shapes.add_picture(img_details, Inches(5), Inches(1.5), height=Inches(5))

    # Slide 5: Historical Tracking
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Feature: Secure Historical Tracking"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.text = "All scans are stored in a Supabase PostgreSQL database for persistent analysis and compliance."
    
    if os.path.exists(img_history):
        slide.shapes.add_picture(img_history, Inches(0.5), Inches(2.5), width=Inches(9))

    # Slide 6: Security Assistant
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Feature: Interactive Security Chat"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.text = "Users can ask follow-up questions to our integrated AI Security Assistant to learn more about cybersecurity."
    
    if os.path.exists(img_chat):
        slide.shapes.add_picture(img_chat, Inches(0.5), Inches(2.5), width=Inches(9))

    # Slide 7: OpenAI Codex - The Ultimate Catalyst
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Special Thanks to OpenAI Codex"
    
    content = slide.placeholders[1]
    content.text = "How Codex powered our development journey:\n" \
                   "1. Accelerated engineering (10x faster development).\n" \
                   "2. Orchestrated the entire Next.js & FastAPI Monorepo.\n" \
                   "3. Handled complex Scikit-learn TF-IDF logic.\n" \
                   "4. Implemented advanced Framer Motion UI/UX automatically.\n" \
                   "Codex isn't just a copilot; it acted as a senior architect!"
                   
    # Save the presentation
    output_path = r"C:\Users\user\Desktop\Spam_Email_Detector\PhishGuard_AI_Hackathon_Pitch.pptx"
    prs.save(output_path)
    print(f"Presentation created successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
